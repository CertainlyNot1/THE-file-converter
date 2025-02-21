from flask import Flask, render_template, request, redirect, url_for, send_file, jsonify
import os
from werkzeug.utils import secure_filename
from moviepy.video.io.VideoFileClip import VideoFileClip

from pdf2docx import Converter
from PIL import Image

import fitz
import pytz
from datetime import datetime
from docx import Document

from pptx import Presentation
from fpdf import FPDF
from bs4 import BeautifulSoup
import requests
from functools import lru_cache



app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['CONVERTED_FOLDER'] = 'converted'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok = True)
os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok = True)

@app.route('/')
def coweb():
    rates = get_ratedd()
    return render_template('index.html', rates = rates)


@app.route('/vid_to_audio', methods=['GET', 'POST'])
def vid_to_audio():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.rsplit('.',1)[0] +'.mp3')
            video = VideoFileClip(filepath)
            video.audio.write_audiofile(output_path)

        return send_file(output_path, as_attachment=True)
    return render_template('video_to_audio.html')

@app.route('/pdf_to_word', methods=['GET', 'POST'])
def pdf_to_word():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            output_filename = filename.rsplit('.',1)[0]+'docx'
            output_path = os.path.join(app.config['CONVERTED_FOLDER'],output_filename)
            doc = Document()
            pdf_doc = fitz.open(filepath)
            for i in pdf_doc:
                text = i.get_text('text')
                doc.add_paragraph(text)
            doc.save(output_path)

            return send_file(output_path, as_attachment=True)
    return render_template('pdf2wrd.html')

@app.route('/image_conversion', methods=['GET', 'POST'])
def image_conversion():
    if request.method == 'POST':
        file = request.files['file']
        conversion_type = request.form.get('conversion_type')

        if file and conversion_type:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            output_format = 'png' if conversion_type == 'to_png' else 'jpg'
            output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.split('.', 1)[0] + f'.{output_format}')

            image = Image.open(filepath)
            if output_format == 'jpg':
                if image.mode in ('RGBA','LA'):
                    background = Image.new('RGB',image.size,(255,255,255))
                    image = Image.alpha_composite(background.convert('RGBA'),image.convert('RGBA')).convert('RGB')
            image.save(output_path,'JPEG' if output_format.upper() == 'JPG' else output_format.upper())

            return send_file(output_path, as_attachment=True)
    return render_template('image_conversion.html')

@app.route('/txt_to_py', methods=['GET', 'POST'])
def txt_to_py():
    if request.method == 'POST':
        file = request.files['file']
        conversion_lang = request.form.get('conversion_lang')
        if file and conversion_lang:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            expendables = {
                'python': 'py',
                'c++': 'cpp',
                'c#': 'cs',
                'javascript': 'js',
                'java': 'java',
                'lua': 'lua',
                'HolyC': 'HC'
            }

            output_filename = filename.split('.', 1)[0] + '.' + expendables[conversion_lang]
            output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)

            with open(filepath, 'r', encoding='utf-8') as txt_file, open(output_path, 'w', encoding='utf-8') as py_file:
                py_file.write(txt_file.read())

            return send_file(output_path, as_attachment=True)
    return render_template('txt_to_py.html')

@app.route('/ppt_to_pdf', methods=['GET', 'POST'])
def ppt_to_pdf():
    if request.method == 'POST':
        file = request.files['file']
        conversion_type = request.form.get('conversion_type')
        if file and conversion_type in ['ppt_to_pdf','pdf_to_ppt']:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            output_filename = filename.split('.', 1)[0] + ('.pdf' if conversion_type == 'ppt_to_pdf' else '.pptx')
            output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)

            if conversion_type == "ppt_to_pdf":

                prs = Presentation(filepath)
                pdf = FPDF()
                pdf.set_auto_page_break(auto=True, margin=15)

                for slide in prs.slides:
                    pdf.add_page()
                    text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    pdf.set_font("Arial", size=12)
                    pdf.multi_cell(0, 10, text)

                pdf.output(output_path)
            
            elif conversion_type == "pdf_to_ppt":
                pdf = fitz.open(filepath)
                preston = Presentation()
                for page in pdf:
                    text = page.get_text('text')
                    slide = preston.slides.add_slide(preston.slide_layouts[5])
                    textbox = slide.shapes.add_textbox(left=10, top=10, width=preston.slide_width, height=preston.slide_height)
                    textbox.text = text
                preston.save(output_path)

            return send_file(output_path, as_attachment=True)
    return render_template('ppt_to_pdf.html')

# @lru_cache(maxsize=32)
def get_rich(in_currency, out_currency):
    url = f"https://www.google.com/finance/quote/{in_currency}-{out_currency}?sa=X&ved=2ahUKEwiwpMTwnoyJAxVIQ_EDHYzuEU0QmY0JegQIARAs"
    content = requests.get(url).text
    soup = BeautifulSoup(content, 'html.parser')
    currency = soup.find("div", class_="YMlKec fxKbKc").get_text()
    return float(currency)

def goldenLeBonBonClatt():
    leUrl = 'https://www.xe.com/currency/'
    LeResponse = requests.get(leUrl)
    LeSoup = BeautifulSoup(LeResponse.text, 'html.parser')
    LeCurrencies = {}
    for LeI in LeSoup.find_all('span',class_='currencyCode'):
        LeCurrencyCode = LeI.get_text(strip = True).replace('-','')
        LeParent = LeI.find_parent()
        LeCurrName = LeParent.get_text(strip=True).replace(LeCurrencyCode,'').replace('-','').strip()
        LeCurrencies[LeCurrencyCode] = LeCurrName
    return LeCurrencies

# @lru_cache(maxsize=32)
def goldenLeBonBonClatt22():
    leUrl = 'https://www.xe.com/currency/'
    LeResponse = requests.get(leUrl)
    LeSoup = BeautifulSoup(LeResponse.text, 'html.parser')
    LeCurrencies = []
    for LeI in LeSoup.find_all('span',class_='currencyCode'):
        LeCurrencyCode = LeI.get_text(strip = True).replace('-','')
        LeCurrencies.append(LeCurrencyCode)
    return LeCurrencies

@app.route('/currency_converter', methods = ['GET','POST'])
def currency_converter():
    result = None
    currencies = goldenLeBonBonClatt()
    if request.method == 'POST':
        from_currency = request.form.get('from_currency')
        to_currency = request.form.get('to_currency')
        amount = float(request.form.get('amount', 1))

        rate = get_rich(from_currency, to_currency)
        result = round(amount * rate, 2)

    return render_template('currency_converter.html', result=result, currencies=currencies)

def get_ratedd():
    currencies = ['EUR','CAD','AED','UAH','KZT','PLN']
    rates = {}
    for currency in currencies:
        try:
            
            rates[currency] = get_rich('USD',currency)
        except Exception as e:
            rates[currency] = 'N/A'

    return rates

def timezones():
    timezones = pytz.all_timezones
    time_data = {}

    for timezone in timezones:
        tz = pytz.timezone(timezone)
        current_time = datetime.now(tz).strftime('%Y-%m-%d %H:%M:%S')
        time_data[timezone] = current_time

    return time_data

@app.route('/GetTimeUhH')
def get_time():
    return jsonify(timezones())

@app.route('/GetRates')
def get_rates():
    rates = get_ratedd()  # Fetch latest rates
    return jsonify(rates)



if __name__ == "__main__":
    app.run(debug=True)
